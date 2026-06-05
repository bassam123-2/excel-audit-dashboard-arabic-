from __future__ import annotations

from email.headerregistry import Address
from email.message import EmailMessage
from email.utils import parseaddr
from io import BytesIO
from pathlib import Path
import base64
import json
from functools import lru_cache
import hashlib
import os
import re
import smtplib
import ssl
import sys
import unicodedata
import warnings

# Excel من القالب يحتوي أحياناً على Data Validation؛ openpyxl يزيلها ويطبع تحذيراً مزعجاً.
warnings.filterwarnings(
    "ignore",
    message="Data Validation extension is not supported",
    category=UserWarning,
    module="openpyxl.worksheet._reader",
)

import pandas as pd
from openpyxl import load_workbook

# شفاف 1×1 بكسل — يُستخدم عند عدم وجود ملف شعار
_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


def _resource_root() -> Path:
    if getattr(sys, "frozen", False) and hasattr(sys, "_MEIPASS"):
        return Path(sys._MEIPASS).resolve()
    return Path(__file__).resolve().parent


def _writable_app_dir() -> Path:
    if getattr(sys, "frozen", False):
        base = Path(os.environ.get("LOCALAPPDATA", str(Path.home()))) / "ExcelArabicDashboard"
        base.mkdir(parents=True, exist_ok=True)
        return base
    return Path(__file__).resolve().parent


def _runtime_app_dir() -> Path:
    """Directory beside the running app executable/script."""
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent


def _brand_logo_search_roots() -> list[Path]:
    roots: list[Path] = []
    for base in (_runtime_app_dir(), _resource_root()):
        if base not in roots:
            roots.append(base)
    return roots


def _pick_brand_logo_code(subsidiary_values: list[str], holding_values: list[str] | None = None) -> str | None:
    """Case-sensitive brand codes; الشركة التابعة then الشركة القابضة."""
    holding_values = holding_values or []
    for values in (subsidiary_values, holding_values):
        for raw in values:
            code = str(raw).strip()
            if code in BRAND_LOGO_CODES:
                return code
    return None


def _brand_logo_path(code: str) -> Path | None:
    for root in _brand_logo_search_roots():
        logos_dir = root / "assets" / "logos"
        for ext in _BRAND_LOGO_EXTENSIONS:
            candidate = logos_dir / f"{code}{ext}"
            if candidate.is_file():
                return candidate.resolve()
    return None


def _brand_logo_mime(path: Path) -> str:
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
    }.get(path.suffix.lower(), "image/png")


def _brand_logo_data_url(code: str | None) -> str:
    if not code:
        return ""
    path = _brand_logo_path(code)
    if path is None:
        return ""
    b64 = base64.standard_b64encode(path.read_bytes()).decode("ascii")
    return f"data:{_brand_logo_mime(path)};base64,{b64}"


def _resolve_logo_path() -> Path | None:
    env = os.environ.get("EXCEL_ARABIC_LOGO", "").strip()
    if env:
        p = Path(env)
        if p.is_file():
            return p.resolve()
    # External override beside the exe/script (no rebuild needed for logo changes).
    runtime_root = _runtime_app_dir()
    for rel in ("logo.png", "brand_logo.png", "logo.jpg", "logo.jpeg", "logo.webp"):
        candidate = runtime_root / rel
        if candidate.is_file():
            return candidate.resolve()
    root = _resource_root()
    for rel in ("assets/logo.png", "assets/brand_logo.png", "static/logo.png"):
        candidate = root / rel
        if candidate.is_file():
            return candidate.resolve()
    return None


LOGO_PATH = _resolve_logo_path()

BRAND_LOGO_CODES = frozenset({"nat", "aum", "saco", "autostar", "btc"})
_BRAND_LOGO_EXTENSIONS = (".png", ".jpg", ".jpeg", ".webp")


UPLOADS_DIR = _writable_app_dir() / "uploads"
UPLOADS_DIR.mkdir(exist_ok=True)


def _load_local_env_file() -> dict[str, str]:
    """Load KEY=VALUE config for desktop/exe runtime."""
    candidates = [
        _runtime_app_dir() / ".env",
        _runtime_app_dir() / "smtp.env",
        _writable_app_dir() / ".env",
        _writable_app_dir() / "smtp.env",
        Path(__file__).resolve().parent / ".env",
        Path(__file__).resolve().parent / "smtp.env",
    ]
    values: dict[str, str] = {}
    for p in candidates:
        if not p.is_file():
            continue
        try:
            for raw in p.read_text(encoding="utf-8").splitlines():
                line = raw.strip()
                if not line or line.startswith("#") or "=" not in line:
                    continue
                k, v = line.split("=", 1)
                key = k.strip()
                if not key:
                    continue
                val = v.strip().strip('"').strip("'")
                values.setdefault(key, val)
        except Exception:
            continue
    return values


_LOCAL_ENV = _load_local_env_file()


def _cfg(name: str, default: str = "") -> str:
    """OS env wins; fallback to local .env/smtp.env."""
    return os.environ.get(name, "").strip() or _LOCAL_ENV.get(name, default).strip()



ARABIC_SHEET_NAME = "سجل الالتزام الموحد"
COL_INHERENT_RISK = "تصنيف المخاطر الكامنة"
COL_RESIDUAL_RISK = "تصنيف المخاطر المتبقية"
COL_STATUS = "الحالة"
COL_TARGET_DATE = "تاريخ التصحيح المستهدف"
COL_YEAR = "السنوات"
COL_DEPARTMENT = "الإدارة المسؤولة"
COL_LEGISLATOR = "المشرع"
COL_SYSTEM_NAME = "اسم النظام"
COL_AUTHORITY = "الهيئة التابعة"
COL_REGULATION = "اللائحة"
COL_LEGAL_TEXT = "النص النظامي"
COL_COMPLIANCE_STATUS = "حالة الالتزام"
COL_CONTROL_CATEGORY = "فئة الضوابط الرقابية"
COL_TASK_OWNER = "مالك المهمة / مالك الإجراء"
COL_RESPONSIBLE_PERSON = "الشخص المسؤول"
COL_CORRECTIVE_PLAN = "الخطة التصحيحية"
COL_MODIFIED_DATE = "تاريخ التصحيح المعدل"
COL_MANAGEMENT_NOTES = "ملاحظات الإدارة.1"
COL_COMPLIANCE_NOTES = "ملاحظات الإلتزام"
COL_HOLDING_COMPANY = "الشركة القابضة"
COL_SUBSIDIARY_COMPANY = "الشركة التابعة"
COL_EMAIL = "email"
YEAR_COL = "السنة"
BLANK_LABEL = "(blank)"
# موضوع بسيط — بعض الخوادم ترفض موضوعاً فارغاً
LEGAL_EMAIL_SUBJECT_PLACEHOLDER = "."
EXTRACTED_IMAGES_DIR = UPLOADS_DIR / "extracted_images"
EXTRACTED_IMAGES_DIR.mkdir(exist_ok=True)


def _normalize_value(value: object) -> str:
    if pd.isna(value):
        return BLANK_LABEL
    text = str(value).strip()
    return text if text else BLANK_LABEL


def _extract_valid_year(value: object) -> str | None:
    """Return a clean 4-digit year if value looks like a real year."""
    if pd.isna(value):
        return None

    if isinstance(value, pd.Timestamp):
        year = int(value.year)
        return str(year) if 1900 <= year <= 2100 else None

    text = str(value).strip()
    if not text:
        return None

    # Numeric forms like 2026 or 2026.0
    try:
        number = float(text)
        if number.is_integer():
            year = int(number)
            if 1900 <= year <= 2100:
                return str(year)
    except ValueError:
        pass

    # Text containing a 4-digit year.
    match = re.search(r"(19|20)\d{2}", text)
    if match:
        return match.group(0)
    return None


def _parse_excel_date_value(value: object):
    """Parse Excel date cells; treat 0/empty placeholders as missing (not 1970)."""
    if pd.isna(value):
        return pd.NaT

    if isinstance(value, pd.Timestamp):
        return value if 1900 <= value.year <= 2100 else pd.NaT

    if isinstance(value, (int, float)) and not isinstance(value, bool):
        number = float(value)
        if number <= 0:
            return pd.NaT
        parsed = pd.to_datetime(number, unit="D", origin="1899-12-30", errors="coerce")
        if pd.notna(parsed) and 1900 <= parsed.year <= 2100:
            return parsed
        return pd.NaT

    text = str(value).strip()
    if not text or text in {"0", "0.0"}:
        return pd.NaT

    parsed = pd.to_datetime(value, errors="coerce")
    if pd.isna(parsed):
        return pd.NaT
    return parsed if 1900 <= parsed.year <= 2100 else pd.NaT


def _parse_excel_date_series(series: pd.Series) -> pd.Series:
    return series.map(_parse_excel_date_value)


def _norm_col_name(name: str) -> str:
    t = unicodedata.normalize("NFKC", str(name).strip())
    return t.lstrip("\ufeff\u200e\u200f").strip()


def _find_optional_email_column(columns: list[str]) -> str | None:
    """عمود البريد: email / E-mail / إلخ (لا يُلزم وجوده)."""
    for col in columns:
        low = _norm_col_name(col).lower().replace("\u00a0", " ").strip()
        compact = low.replace(" ", "").replace("-", "").replace("_", "")
        # مطابقة صارمة لتفادي أعمدة عربية طويلة تحتوي حرفياً على «email»
        if re.fullmatch(r"email(\.\d+)?", compact, flags=re.IGNORECASE):
            return col
        if re.fullmatch(r"e[-_]mail(\.\d+)?", compact, flags=re.IGNORECASE):
            return col
        if "البريد" in low and ("إلكتروني" in low or "الكتروني" in low):
            return col
    return None


def _find_physical_column(columns: list[str], preferred: str, required_tokens: list[str]) -> str:
    normalized_map = {_norm_col_name(c): c for c in columns}
    preferred_norm = _norm_col_name(preferred)
    if preferred_norm in normalized_map:
        return normalized_map[preferred_norm]

    candidates: list[str] = []
    for col in columns:
        norm = _norm_col_name(col)
        if all(token in norm for token in required_tokens):
            candidates.append(col)

    if not candidates:
        raise KeyError(preferred)

    if ".1" in preferred_norm:
        for col in candidates:
            if ".1" in _norm_col_name(col):
                return col

    return candidates[0]


def _find_physical_column_by_rank(
    columns: list[str],
    preferred: str,
    required_tokens: list[str],
    rank_token_sets: list[list[str]],
) -> str:
    """Like _find_physical_column, but pick the best candidate using rank_token_sets order."""
    normalized_map = {_norm_col_name(c): c for c in columns}
    preferred_norm = _norm_col_name(preferred)
    if preferred_norm in normalized_map:
        return normalized_map[preferred_norm]

    candidates: list[str] = []
    for col in columns:
        norm = _norm_col_name(col)
        if all(token in norm for token in required_tokens):
            candidates.append(col)

    if not candidates:
        raise KeyError(preferred)

    if ".1" in preferred_norm:
        for col in candidates:
            if ".1" in _norm_col_name(col):
                return col

    for token_set in rank_token_sets:
        for col in candidates:
            norm = _norm_col_name(col)
            if all(token in norm for token in token_set):
                return col

    return candidates[0]


def _find_system_name_physical_column(columns: list[str]) -> str:
    """اسم النظام — أو «النظام» فقط في بعض قوالب سجل الالتزام."""
    normalized_map = {_norm_col_name(c): c for c in columns}
    preferred_norm = _norm_col_name(COL_SYSTEM_NAME)
    if preferred_norm in normalized_map:
        return normalized_map[preferred_norm]
    try:
        return _find_physical_column(columns, COL_SYSTEM_NAME, ["اسم", "النظام"])
    except KeyError:
        pass
    system_only = _norm_col_name("النظام")
    if system_only in normalized_map:
        return normalized_map[system_only]
    raise KeyError(COL_SYSTEM_NAME)


def _find_inherent_risk_physical_column(columns: list[str]) -> str:
    return _find_physical_column_by_rank(
        columns,
        COL_INHERENT_RISK,
        ["مخاطر", "كامنة"],
        [
            ["تصنيف", "كامنة"],
            ["مستوى", "كامنة"],
            ["درجة", "كامنة"],
        ],
    )


def _find_compliance_status_physical_column(columns: list[str]) -> str:
    normalized_map = {_norm_col_name(c): c for c in columns}
    preferred_norm = _norm_col_name(COL_COMPLIANCE_STATUS)
    if preferred_norm in normalized_map:
        return normalized_map[preferred_norm]

    candidates = [
        col
        for col in columns
        if all(token in _norm_col_name(col) for token in ["حالة", "التزام"])
    ]
    if not candidates:
        raise KeyError(COL_COMPLIANCE_STATUS)

    for col in candidates:
        if "إدارة الالتزام" in _norm_col_name(col):
            return col
    return candidates[0]


def _find_holding_company_physical_column(columns: list[str]) -> str:
    normalized_map = {_norm_col_name(c): c for c in columns}
    preferred_norm = _norm_col_name(COL_HOLDING_COMPANY)
    if preferred_norm in normalized_map:
        return normalized_map[preferred_norm]
    for col in columns:
        n = _norm_col_name(col)
        if "الشركة" in n and "قابضة" in n:
            return col
    raise KeyError(COL_HOLDING_COMPANY)


def _find_subsidiary_company_physical_column(columns: list[str]) -> str:
    """الشركة التابعة — لا تُخلط مع «الهيئة التابعة»."""
    normalized_map = {_norm_col_name(c): c for c in columns}
    preferred_norm = _norm_col_name(COL_SUBSIDIARY_COMPANY)
    if preferred_norm in normalized_map:
        return normalized_map[preferred_norm]
    for col in columns:
        n = _norm_col_name(col)
        if "الشركة" in n and "تابعة" in n and "قابضة" not in n:
            return col
    raise KeyError(COL_SUBSIDIARY_COMPANY)


def _find_corrective_plan_physical_column(columns: list[str]) -> str:
    for tokens in (["الخطة", "تصحيح"], ["الإجراء", "تصحيح"]):
        try:
            return _find_physical_column(columns, COL_CORRECTIVE_PLAN, tokens)
        except KeyError:
            continue
    raise KeyError(COL_CORRECTIVE_PLAN)


def _find_legal_text_physical_column(columns: list[str]) -> str:
    """النص النظامي، أو بديله الشائع «النص بالكامل» في بعض القوالب."""
    try:
        return _find_physical_column(columns, COL_LEGAL_TEXT, ["النص", "نظام"])
    except KeyError:
        pass
    try:
        return _find_physical_column(columns, COL_LEGAL_TEXT, ["النص", "كامل"])
    except KeyError:
        pass
    for col in columns:
        n = _norm_col_name(col)
        if "النص" in n and "كامل" in n:
            return col
    raise KeyError(COL_LEGAL_TEXT)


def _rename_columns_to_canonical(df: pd.DataFrame) -> pd.DataFrame:
    columns = [str(c).strip() for c in df.columns]
    mapping = {
        _find_inherent_risk_physical_column(columns): COL_INHERENT_RISK,
        _find_physical_column(columns, COL_RESIDUAL_RISK, ["مخاطر", "متبق"]): COL_RESIDUAL_RISK,
        _find_physical_column(columns, COL_STATUS, ["الحالة"]): COL_STATUS,
        _find_physical_column(columns, COL_DEPARTMENT, ["الإدارة", "مسؤول"]): COL_DEPARTMENT,
        _find_physical_column(columns, COL_LEGISLATOR, ["المشرع"]): COL_LEGISLATOR,
        _find_system_name_physical_column(columns): COL_SYSTEM_NAME,
        _find_physical_column(columns, COL_AUTHORITY, ["الهيئة", "تابعة"]): COL_AUTHORITY,
        _find_physical_column(columns, COL_REGULATION, ["اللائحة"]): COL_REGULATION,
        _find_legal_text_physical_column(columns): COL_LEGAL_TEXT,
        _find_compliance_status_physical_column(columns): COL_COMPLIANCE_STATUS,
        _find_physical_column(columns, COL_CONTROL_CATEGORY, ["فئة", "ضوابط"]): COL_CONTROL_CATEGORY,
    }

    if COL_YEAR in columns or _norm_col_name(COL_YEAR) in {_norm_col_name(c) for c in columns}:
        mapping[_find_physical_column(columns, COL_YEAR, ["السنوات"])] = COL_YEAR

    if COL_TARGET_DATE in columns or any("التصحيح" in _norm_col_name(c) and "المستهدف" in _norm_col_name(c) for c in columns):
        mapping[_find_physical_column(columns, COL_TARGET_DATE, ["التصحيح", "المستهدف"])] = COL_TARGET_DATE

    email_physical = _find_optional_email_column(columns)
    if email_physical and email_physical not in mapping:
        mapping[email_physical] = COL_EMAIL

    # Optional / detail columns (may exist)
    optional_specs = [
        (COL_TASK_OWNER, ["مالك", "مهمة"]),
        (COL_RESPONSIBLE_PERSON, ["الشخص", "مسؤول"]),
        (COL_MODIFIED_DATE, ["التصحيح", "المعدل"]),
        (COL_MANAGEMENT_NOTES, ["ملاحظات", "الإدارة"]),
        (COL_COMPLIANCE_NOTES, ["ملاحظات", "الإلتزام"]),
    ]
    for preferred, tokens in optional_specs:
        try:
            physical = _find_physical_column(columns, preferred, tokens)
        except KeyError:
            continue
        if physical in mapping:
            continue
        mapping[physical] = preferred

    for finder, canonical in (
        (_find_holding_company_physical_column, COL_HOLDING_COMPANY),
        (_find_subsidiary_company_physical_column, COL_SUBSIDIARY_COMPANY),
    ):
        try:
            physical = finder(columns)
        except KeyError:
            continue
        if physical in mapping:
            continue
        mapping[physical] = canonical

    try:
        corrective_physical = _find_corrective_plan_physical_column(columns)
        if corrective_physical not in mapping:
            mapping[corrective_physical] = COL_CORRECTIVE_PLAN
    except KeyError:
        pass

    out = df.rename(columns=mapping)
    out.columns = [_norm_col_name(c) for c in out.columns]
    return out


def _workbook_primary_sheet(wb: object) -> object | None:
    if ARABIC_SHEET_NAME in wb.sheetnames:
        return wb[ARABIC_SHEET_NAME]
    if wb.sheetnames:
        return wb[wb.sheetnames[0]]
    return None


def _scan_email_column_index_uncached(path: Path) -> int | None:
    """عمود Excel (1-based) الذي عنوانه email في أحد الصفوف الأولى."""
    wb = None
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        ws = _workbook_primary_sheet(wb)
        if ws is None:
            return None
        max_r = 12
        max_c = 400
        for row in ws.iter_rows(min_row=1, max_row=max_r, min_col=1, max_col=max_c, values_only=True):
            for c_idx, val in enumerate(row, start=1):
                if val is None:
                    continue
                s = _norm_col_name(str(val)).lower().replace("\u00a0", " ")
                if re.fullmatch(r"email|e-mail|e_mail|الإيميل|الايميل", s):
                    return int(c_idx)
        return None
    except Exception:
        return None
    finally:
        if wb is not None:
            wb.close()


@lru_cache(maxsize=32)
def _find_email_excel_column_1based_cached(path_str: str, token: str) -> int | None:
    return _scan_email_column_index_uncached(Path(path_str))


def _find_email_excel_column_1based(path: Path) -> int | None:
    try:
        return _find_email_excel_column_1based_cached(str(path.resolve()), _file_cache_token(path))
    except OSError:
        return _scan_email_column_index_uncached(path)


def _read_email_at_excel_row(path: Path, excel_row_1based: int) -> str:
    """قراءة خلية البريد من الإكسل مباشرة (نفس رقم صف الصف في الملف) — يتجاوز اختلاف أعمدة pandas."""
    if excel_row_1based < 1:
        return ""
    ec = _find_email_excel_column_1based(path)
    if not ec:
        return ""

    def _read_ro() -> str:
        wb = None
        try:
            wb = load_workbook(path, read_only=True, data_only=True)
            ws = _workbook_primary_sheet(wb)
            if ws is None:
                return ""
            for row in ws.iter_rows(
                min_row=excel_row_1based,
                max_row=excel_row_1based,
                min_col=ec,
                max_col=ec,
                values_only=True,
            ):
                if not row:
                    return ""
                v = row[0]
                if v is None:
                    return ""
                return str(v).strip()
            return ""
        except Exception:
            return ""
        finally:
            if wb is not None:
                wb.close()

    def _read_rw() -> str:
        try:
            wb = load_workbook(path, read_only=False, data_only=True)
            ws = _workbook_primary_sheet(wb)
            if ws is None:
                return ""
            c = ws.cell(row=excel_row_1based, column=ec)
            v = c.value
            if v is None:
                return ""
            return str(v).strip()
        except Exception:
            return ""

    out = _read_ro()
    return out if out else _read_rw()


def _nonempty_email_mask(series: pd.Series) -> pd.Series:
    s = series.astype(str).str.strip()
    return series.notna() & s.ne("") & s.str.contains(r"@", regex=True, na=False)


def _sync_email_column_from_excel(path: Path, df: pd.DataFrame) -> pd.DataFrame:
    """
    إذا كان عنوان «email» على صف غير صف header=3 يظهر عند pandas كـ Unnamed ولا يُربط.
    نمسح الصفوف الأولى بالورقة وننسخ عمود البيانات من نفس فهرس عمود Excel.
    """
    ecol = _find_email_excel_column_1based(path)
    if ecol is None:
        return df
    idx = int(ecol) - 1
    if idx < 0 or idx >= len(df.columns):
        return df
    scanned = df.iloc[:, idx].copy()
    out = df.copy()
    if COL_EMAIL not in out.columns:
        out[COL_EMAIL] = scanned
        return out
    if _nonempty_email_mask(scanned).sum() > _nonempty_email_mask(out[COL_EMAIL]).sum():
        out[COL_EMAIL] = scanned
    return out


_EMAIL_ADDR_RE = re.compile(r"^[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}$")


def _normalize_recipient_email(raw: object) -> str | None:
    if raw is None or (isinstance(raw, float) and pd.isna(raw)):
        return None
    t = str(raw).strip()
    if not t or t == BLANK_LABEL:
        return None
    if _EMAIL_ADDR_RE.match(t):
        return t
    m = re.search(r"[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}", t)
    return m.group(0) if m else None


def _smtp_config() -> dict[str, str | int | bool] | None:
    host = _cfg("EXCEL_ARABIC_SMTP_HOST")
    user = _cfg("EXCEL_ARABIC_SMTP_USER")
    password = _cfg("EXCEL_ARABIC_SMTP_PASSWORD")
    if not host or not user or not password:
        return None
    port = int(_cfg("EXCEL_ARABIC_SMTP_PORT", "587"))
    use_tls = _cfg("EXCEL_ARABIC_SMTP_USE_TLS", "1").lower() not in ("0", "false", "no")
    from_addr = _cfg("EXCEL_ARABIC_SMTP_FROM") or user
    return {"host": host, "port": port, "user": user, "password": password, "from": from_addr, "use_tls": use_tls}


def _smtp_from_address(cfg: dict[str, str | int | bool]) -> Address:
    """حقل From مع اسم العرض العربي (يظهر لدى المستلم بدل عنوان البريد فقط)."""
    display = _cfg("EXCEL_ARABIC_SMTP_DISPLAY_NAME", "ادارة الالتزام") or "ادارة الالتزام"
    _old_name, addr = parseaddr(str(cfg["from"]))
    if not addr:
        addr = str(cfg["from"]).strip()
    return Address(display_name=display, addr_spec=addr)


def _send_smtp_message(to_addr: str, subject: str, body: str) -> None:
    cfg = _smtp_config()
    if not cfg:
        raise RuntimeError(
            "إعداد SMTP غير مكتمل. عيّن EXCEL_ARABIC_SMTP_HOST و EXCEL_ARABIC_SMTP_USER و EXCEL_ARABIC_SMTP_PASSWORD "
            "(واختياري: EXCEL_ARABIC_SMTP_FROM، EXCEL_ARABIC_SMTP_PORT، EXCEL_ARABIC_SMTP_USE_TLS)."
        )
    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"] = _smtp_from_address(cfg)
    msg["To"] = to_addr
    msg.set_content(body)
    if cfg["use_tls"]:
        ctx = ssl.create_default_context()
        with smtplib.SMTP(str(cfg["host"]), int(cfg["port"]), timeout=45) as smtp:
            smtp.starttls(context=ctx)
            smtp.login(str(cfg["user"]), str(cfg["password"]))
            smtp.send_message(msg)
    else:
        with smtplib.SMTP(str(cfg["host"]), int(cfg["port"]), timeout=45) as smtp:
            smtp.login(str(cfg["user"]), str(cfg["password"]))
            smtp.send_message(msg)


def _file_cache_token(path: Path) -> str:
    st = path.stat()
    return f"{path.resolve()}::{st.st_mtime_ns}::{st.st_size}"


def _load_dashboard_df_uncached(path: Path) -> pd.DataFrame:
    df = pd.read_excel(path, sheet_name=ARABIC_SHEET_NAME, header=3)
    df = _rename_columns_to_canonical(df)
    df = _sync_email_column_from_excel(path, df)

    required_cols = [
        COL_INHERENT_RISK,
        COL_RESIDUAL_RISK,
        COL_STATUS,
        COL_DEPARTMENT,
        COL_LEGISLATOR,
        COL_SYSTEM_NAME,
        COL_AUTHORITY,
        COL_REGULATION,
        COL_LEGAL_TEXT,
        COL_COMPLIANCE_STATUS,
        COL_CONTROL_CATEGORY,
    ]
    missing = [col for col in required_cols if col not in df.columns]
    if missing:
        raise ValueError(f"Missing columns: {', '.join(missing)}")

    filtered = df[required_cols].copy()
    target_dates = _parse_excel_date_series(df[COL_TARGET_DATE]) if COL_TARGET_DATE in df.columns else None

    if COL_YEAR in df.columns:
        year_from_col = df[COL_YEAR].map(_extract_valid_year)
        if target_dates is not None:
            year_from_date = target_dates.dt.year.where(target_dates.notna()).astype("Int64").astype(str)
            year_from_date = year_from_date.where(target_dates.notna(), None)
            filtered[YEAR_COL] = year_from_col.fillna(year_from_date).fillna(BLANK_LABEL)
        else:
            filtered[YEAR_COL] = year_from_col.fillna(BLANK_LABEL)
    elif target_dates is not None:
        filtered[YEAR_COL] = target_dates.dt.year.astype("Int64").astype(str)
        filtered.loc[target_dates.isna(), YEAR_COL] = BLANK_LABEL
    else:
        filtered[YEAR_COL] = BLANK_LABEL

    for col in [
        COL_INHERENT_RISK,
        COL_RESIDUAL_RISK,
        COL_STATUS,
        COL_DEPARTMENT,
        COL_LEGISLATOR,
        COL_SYSTEM_NAME,
        COL_AUTHORITY,
        COL_REGULATION,
        COL_LEGAL_TEXT,
        COL_COMPLIANCE_STATUS,
        COL_CONTROL_CATEGORY,
    ]:
        filtered[col] = filtered[col].map(_normalize_value)

    return filtered


def _load_full_df_uncached(path: Path) -> pd.DataFrame:
    df = pd.read_excel(path, sheet_name=ARABIC_SHEET_NAME, header=3)
    df = _rename_columns_to_canonical(df)
    df = _sync_email_column_from_excel(path, df)
    return df


@lru_cache(maxsize=16)
def _load_dashboard_df_cached(path_str: str, token: str) -> pd.DataFrame:
    _ = token
    return _load_dashboard_df_uncached(Path(path_str))


@lru_cache(maxsize=16)
def _load_full_df_cached(path_str: str, token: str) -> pd.DataFrame:
    _ = token
    return _load_full_df_uncached(Path(path_str))


def _load_dashboard_df(path: Path) -> pd.DataFrame:
    return _load_dashboard_df_cached(str(path), _file_cache_token(path))


def _load_full_df(path: Path) -> pd.DataFrame:
    return _load_full_df_cached(str(path), _file_cache_token(path))


FILTER_DIMENSION_COLS = [
    COL_INHERENT_RISK,
    COL_RESIDUAL_RISK,
    COL_STATUS,
    COL_DEPARTMENT,
    COL_LEGISLATOR,
    COL_SYSTEM_NAME,
    COL_AUTHORITY,
    COL_REGULATION,
    COL_LEGAL_TEXT,
    COL_COMPLIANCE_STATUS,
    COL_CONTROL_CATEGORY,
]

OPTIONAL_FILTER_DIMENSION_COLS = [
    COL_HOLDING_COMPANY,
    COL_SUBSIDIARY_COMPANY,
]


def _filter_dimension_cols_in_df(df: pd.DataFrame) -> list[str]:
    cols = [c for c in FILTER_DIMENSION_COLS if c in df.columns]
    cols.extend(c for c in OPTIONAL_FILTER_DIMENSION_COLS if c in df.columns)
    return cols

AGING_TIME_ROWS = [
    ("not_due", "لم يحن بعد"),
    ("lt_6m", "أقل من 6 أشهر"),
    ("lt_1y", "أقل من سنة"),
    ("ge_1y", "أكثر من سنة"),
]

AGING_RISK_COLS = [
    ("very_high", "مرتفع جدا", "#b91c1c"),
    ("high", "مرتفع", "#ea580c"),
    ("medium", "متوسط", "#ca8a04"),
    ("low", "منخفض", "#16a34a"),
    ("very_low", "متدني الانخفاض", "#86efac"),
    ("other", "أخرى", "#94a3b8"),
]


def _add_year_column_to_df(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    target_dates = _parse_excel_date_series(out[COL_TARGET_DATE]) if COL_TARGET_DATE in out.columns else None

    if COL_YEAR in out.columns:
        year_from_col = out[COL_YEAR].map(_extract_valid_year)
        if target_dates is not None:
            year_from_date = target_dates.dt.year.where(target_dates.notna()).astype("Int64").astype(str)
            year_from_date = year_from_date.where(target_dates.notna(), None)
            out[YEAR_COL] = year_from_col.fillna(year_from_date).fillna(BLANK_LABEL)
        else:
            out[YEAR_COL] = year_from_col.fillna(BLANK_LABEL)
    elif target_dates is not None:
        out[YEAR_COL] = target_dates.dt.year.astype("Int64").astype(str)
        out.loc[target_dates.isna(), YEAR_COL] = BLANK_LABEL
    else:
        out[YEAR_COL] = BLANK_LABEL
    return out


def _normalize_filter_dimensions(df: pd.DataFrame) -> pd.DataFrame:
    out = df.copy()
    for col in _filter_dimension_cols_in_df(out):
        out[col] = out[col].map(_normalize_value)
    if YEAR_COL in out.columns:
        out[YEAR_COL] = out[YEAR_COL].map(_normalize_value)
    return out


def _prepare_filtered_full_df(path: Path, selected: dict[str, list[str]]) -> pd.DataFrame:
    df = _load_full_df(path)
    df = _add_year_column_to_df(df)
    df = _normalize_filter_dimensions(df)
    return _apply_filters(df, selected)


AUDIT_PLAN_PREFERRED: list[str] = [
    COL_CORRECTIVE_PLAN,
    COL_TARGET_DATE,
    COL_MODIFIED_DATE,
    COL_TASK_OWNER,
    COL_RESPONSIBLE_PERSON,
    COL_STATUS,
    COL_COMPLIANCE_STATUS,
    COL_CONTROL_CATEGORY,
    COL_DEPARTMENT,
]

AUDIT_PLAN_DATE_COLS = {COL_TARGET_DATE, COL_MODIFIED_DATE}
AUDIT_PLAN_VALUECAP = 80


def _audit_plan_column_list(df: pd.DataFrame) -> list[str]:
    """أعمدة ذات صلة بمتابعة خطة التصحيح/التدقيق، زائد أي عمود اسمه يحتوي «تدقيق»."""
    found: list[str] = []
    seen: set[str] = set()
    for col in AUDIT_PLAN_PREFERRED:
        if col in df.columns and col not in seen:
            found.append(col)
            seen.add(col)
    for col in df.columns:
        if col in seen:
            continue
        if "تدقيق" in _norm_col_name(str(col)):
            found.append(col)
            seen.add(col)
    return found


def _series_labels_for_audit_column(col: str, series: pd.Series) -> pd.Series:
    if col in AUDIT_PLAN_DATE_COLS:
        dt = _parse_excel_date_series(series)
        out = dt.dt.strftime("%Y-%m-%d")
        return out.where(dt.notna(), BLANK_LABEL)
    return series.map(_normalize_value)


def _is_open_status_for_aging(status_text: str) -> bool:
    """يُحسب التقادم فقط لصفوف: مفتوح (تجاوز) أو مفتوح (ضمن)، مع تسامح بسيط بالمسافات."""
    t = _norm_col_name(str(status_text)).replace("  ", " ")
    if "مفتوح" not in t:
        return False
    if "تجاوز" in t and "تاريخ" in t and "التصحيح" in t:
        return True
    if "ضمن" in t and "تاريخ" in t and "التصحيح" in t:
        return True
    return False


def _aging_risk_key(residual_norm: str) -> str | None:
    """تصنيف المخاطر المتبقية مع تسامح لأخطاء الإدخال الشائعة في الملف."""
    if residual_norm == BLANK_LABEL:
        return None
    raw = str(residual_norm).strip()
    if not raw:
        return None

    t = _norm_col_name(raw).replace("\u00a0", " ")
    t = re.sub(r"\s+", " ", t).strip()
    for bad, good in (
        ("مرنفع", "مرتفع"),
        ("مرتفغ", "مرتفع"),
        ("مرتفاع", "مرتفع"),
        ("مرنفغ", "مرتفع"),
        ("مرتفغ جدا", "مرتفع جدا"),
    ):
        t = t.replace(bad, good)

    if "متدني" in t and (
        "انخفاض" in t or "انخفاظ" in t or "انخغاض" in t or "انخفاق" in t
    ):
        return "very_low"

    has_j = "جدا" in t or "جداً" in t or "جدآ" in t
    if has_j and ("مرتفع" in t or "مرفع" in t):
        return "very_high"

    if "متوسط" in t:
        return "medium"

    if "منخفض" in t and "متدني" not in t:
        return "low"

    if "مرتفع" in t:
        return "high"

    if "مرفع" in t:
        return "high"

    return None


def _aging_time_bucket(compare: pd.Timestamp, reference: pd.Timestamp) -> str | None:
    cref = compare.normalize()
    rref = reference.normalize()
    if pd.isna(cref) or pd.isna(rref):
        return None
    if cref > rref:
        return "not_due"
    overdue_days = int((rref - cref).days)
    if overdue_days < 183:
        return "lt_6m"
    if overdue_days < 365:
        return "lt_1y"
    return "ge_1y"


def _format_display_value(value: object) -> str:
    if pd.isna(value):
        return ""
    if isinstance(value, pd.Timestamp):
        return value.strftime("%Y-%m-%d")
    return str(value).strip()


def _pick_best_legal_match(matches: pd.DataFrame) -> pd.Series:
    """Return the most informative row among legal-text matches."""
    if matches.empty:
        raise ValueError("matches is empty")

    work = matches
    if COL_EMAIL in matches.columns:
        has_mail = matches[COL_EMAIL].map(lambda v: _normalize_recipient_email(v) is not None)
        if bool(has_mail.any()):
            work = matches.loc[has_mail]

    score_cols = [
        COL_STATUS,
        COL_RESIDUAL_RISK,
        COL_CONTROL_CATEGORY,
        COL_TARGET_DATE,
        COL_MODIFIED_DATE,
        COL_TASK_OWNER,
        COL_RESPONSIBLE_PERSON,
        COL_CORRECTIVE_PLAN,
        COL_MANAGEMENT_NOTES,
        COL_COMPLIANCE_NOTES,
    ]
    if COL_EMAIL in work.columns:
        score_cols = [COL_EMAIL, *score_cols]
    available = [c for c in score_cols if c in work.columns]
    if not available:
        return work.iloc[0]

    scored = work.copy()
    score = pd.Series(0, index=scored.index, dtype="int64")
    for c in available:
        if c == COL_EMAIL:
            score = score + scored[c].map(lambda v: 1 if _normalize_recipient_email(v) else 0)
        else:
            score = score + scored[c].map(lambda v: 0 if _normalize_value(v) == BLANK_LABEL else 1)
    best_idx = int(score.idxmax())
    return work.loc[best_idx]


def _excel_row_for_df_index(df_index: int) -> int:
    # header=3 => headers on Excel row 4 (1-based), first data row is Excel row 5.
    return int(df_index) + 5


def _image_signature(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _extract_images_index_uncached(path: Path) -> dict[int, list[str]]:
    """Map Excel row (1-based) -> extracted image URLs for that row."""
    wb = load_workbook(path, read_only=False, data_only=True)
    if ARABIC_SHEET_NAME not in wb.sheetnames:
        return {}
    ws = wb[ARABIC_SHEET_NAME]

    out: dict[int, list[str]] = {}
    if not getattr(ws, "_images", None):
        return {}

    for image in ws._images:  # type: ignore[attr-defined]
        anchor = getattr(image, "anchor", None)
        from_cell = getattr(anchor, "_from", None) if anchor is not None else None
        if from_cell is None:
            continue

        row0 = int(getattr(from_cell, "row", 0))
        img_row_1based = row0 + 1

        data = image._data()  # type: ignore[attr-defined]
        ext = Path(str(getattr(image, "path", "")) or "").suffix.lower()
        if ext not in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
            ext = ".png"

        sig = _image_signature(data)
        filename = f"{sig}{ext}"
        out_path = EXTRACTED_IMAGES_DIR / filename
        if not out_path.exists():
            out_path.write_bytes(data)

        out.setdefault(img_row_1based, []).append(f"/extracted/{filename}")

    return out


@lru_cache(maxsize=16)
def _extract_images_index_cached(path_str: str, token: str) -> dict[int, list[str]]:
    _ = token
    return _extract_images_index_uncached(Path(path_str))


def _extract_images_for_excel_row(path: Path, excel_row: int) -> list[str]:
    """Return URL paths (relative to app) for images anchored on the given Excel row."""
    index = _extract_images_index_cached(str(path), _file_cache_token(path))
    return index.get(int(excel_row), [])


def _image_url_to_path(url: str) -> Path | None:
    u = (url or "").strip()
    if not u:
        return None
    if u.startswith("/extracted/"):
        name = u.split("/extracted/", 1)[1].strip("/")
        if not name or ".." in name.replace("\\", "/"):
            return None
        p = EXTRACTED_IMAGES_DIR / name
        return p if p.is_file() else None
    return None


def _collect_legal_text_record(path: Path, legal_text: str, *, include_images: bool) -> dict[str, object]:
    """Build legal-text detail payload from the uploaded Excel file."""
    df = _load_full_df(path)
    if COL_LEGAL_TEXT not in df.columns:
        raise ValueError("Legal text column missing")

    normalized = df[COL_LEGAL_TEXT].map(_normalize_value)
    lookup_key = _normalize_value(legal_text)
    matches = df[normalized == lookup_key]
    if matches.empty:
        raise LookupError("Not found")

    row = _pick_best_legal_match(matches)
    idx = int(row.name)
    excel_row = _excel_row_for_df_index(idx)

    year_value = ""
    if COL_YEAR in df.columns:
        year_value = _extract_valid_year(row[COL_YEAR]) or ""
    if not year_value and COL_TARGET_DATE in df.columns:
        parsed = _parse_excel_date_value(row[COL_TARGET_DATE])
        if pd.notna(parsed):
            year_value = str(int(parsed.year))

    detail_fields = [
        (COL_STATUS, "الحالة"),
        (COL_RESIDUAL_RISK, COL_RESIDUAL_RISK),
        (COL_CONTROL_CATEGORY, "فئة الضوابط الرقابية"),
        (COL_TARGET_DATE, "تاريخ التصحيح المستهدف"),
        (COL_MODIFIED_DATE, "تاريخ التصحيح المعدل"),
        (COL_TASK_OWNER, "مالك المهمة / مالك الإجراء"),
        (COL_RESPONSIBLE_PERSON, "الشخص المسؤول"),
        (COL_CORRECTIVE_PLAN, "الخطة التصحيحية"),
        (COL_MANAGEMENT_NOTES, "ملاحظات الإدارة"),
        (COL_COMPLIANCE_NOTES, "ملاحظات الإلتزام"),
        (COL_HOLDING_COMPANY, COL_HOLDING_COMPANY),
        (COL_SUBSIDIARY_COMPANY, COL_SUBSIDIARY_COMPANY),
    ]

    fields: list[dict[str, str]] = []
    for col, label in detail_fields:
        if col not in df.columns:
            continue
        val = _format_display_value(row[col])
        fields.append({"label": label, "value": val if val else "—"})

    if year_value:
        fields.insert(3, {"label": "السنوات", "value": year_value})

    cell_email = _read_email_at_excel_row(path, excel_row).strip()
    em_disp = cell_email
    if not em_disp and COL_EMAIL in df.columns:
        raw_m = row[COL_EMAIL] if COL_EMAIL in row.index else None
        em_disp = (
            _format_display_value(raw_m).strip()
            if raw_m is not None and not pd.isna(raw_m)
            else ""
        )
    if _find_email_excel_column_1based(path) is not None or COL_EMAIL in df.columns:
        fields.insert(0, {"label": "البريد الإلكتروني (email)", "value": em_disp or "—"})

    images: list[str] = []
    if include_images:
        images = _extract_images_for_excel_row(path, excel_row)

    return {
        "legal_text": lookup_key,
        "excel_row": excel_row,
        "picked_row_index": idx,
        "recipient_email": em_disp or "",
        "fields": fields,
        "images": images,
    }


def _pptx_rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _pptx_set_rtl_paragraph(paragraph, *, bold: bool = False, size_pt: int = 12, color=None) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    paragraph.alignment = PP_ALIGN.RIGHT
    if not paragraph.runs:
        paragraph.add_run()
    run = paragraph.runs[0]
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _pptx_export_fields(fields: list[dict[str, str]]) -> list[dict[str, str]]:
    """Fields for PPT export: exclude email columns."""
    out: list[dict[str, str]] = []
    for field in fields:
        label = str(field.get("label", "")).strip()
        low = label.lower()
        if "email" in low or "بريد" in label or "البريد" in label:
            continue
        out.append({"label": label, "value": str(field.get("value", "") or "—")})
    return out


def _pptx_legal_text_layout(text: str, width_in: float, max_h_in: float, min_h_in: float = 0.5) -> tuple[float, int]:
    """Pick font size and box height so the full legal text fits (single-slide layout)."""
    t = (text or "—").strip()
    n = len(t)
    if n <= 80:
        return (min(max_h_in, max(min_h_in, 0.55)), 12)
    for pt in (12, 11, 10, 9, 8, 7, 6):
        chars_per_line = max(32, int(width_in * (5.8 - (12 - min(pt, 12)) * 0.35)))
        lines = max(1, (n + chars_per_line - 1) // chars_per_line)
        h = 0.14 + lines * (pt * 0.0172)
        if h <= max_h_in:
            return (max(min_h_in, h), pt)
    return (max_h_in, 6)


def _pptx_fill_cell_compact(cell, text: str, *, header: bool = False, value: bool = False, size_pt: int = 10) -> None:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    cell.text = str(text or "—")
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(2)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.RIGHT
        for run in p.runs:
            run.font.size = Pt(size_pt)
            run.font.bold = header or (not value)
            if header:
                run.font.color.rgb = _pptx_rgb("#ffffff")
            elif value:
                run.font.color.rgb = _pptx_rgb("#111827")
            else:
                run.font.color.rgb = _pptx_rgb("#14532d")
    fill = cell.fill
    fill.solid()
    if header:
        fill.fore_color.rgb = _pptx_rgb("#14532d")
    elif value:
        fill.fore_color.rgb = _pptx_rgb("#ffffff")
    else:
        fill.fore_color.rgb = _pptx_rgb("#ecfdf5")


def _pptx_add_header_bar(slide, prs, margin_x: float, content_w: float, header_label: str) -> float:
    """Green top bar; returns bar height in inches."""
    from pptx.util import Inches

    header_bar_h = 0.42
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(header_bar_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _pptx_rgb("#76b447")
    bar.line.fill.background()

    label_left = margin_x
    if LOGO_PATH and LOGO_PATH.is_file():
        try:
            slide.shapes.add_picture(str(LOGO_PATH), Inches(margin_x), Inches(0.05), height=Inches(0.32))
            label_left = margin_x + 0.55
        except Exception:
            pass

    lbl = slide.shapes.add_textbox(Inches(label_left), Inches(0.08), Inches(content_w - 0.5), Inches(0.3))
    lp = lbl.text_frame.paragraphs[0]
    lp.text = header_label
    _pptx_set_rtl_paragraph(lp, bold=True, size_pt=12, color=_pptx_rgb("#000000"))
    return header_bar_h


def _pptx_add_legal_text_block(
    slide,
    legal_text: str,
    *,
    top_in: float,
    height_in: float,
    margin_x: float,
    content_w: float,
    font_pt: int = 11,
) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    legal_bg = slide.shapes.add_shape(
        1, Inches(margin_x), Inches(top_in), Inches(content_w), Inches(height_in)
    )
    legal_bg.fill.solid()
    legal_bg.fill.fore_color.rgb = _pptx_rgb("#f0fdf4")
    legal_bg.line.color.rgb = _pptx_rgb("#14532d")

    inner_h = max(0.35, height_in - 0.14)
    legal_box = slide.shapes.add_textbox(
        Inches(margin_x + 0.1),
        Inches(top_in + 0.07),
        Inches(content_w - 0.2),
        Inches(inner_h),
    )
    ltf = legal_box.text_frame
    ltf.word_wrap = True
    ltf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    ltf.margin_left = ltf.margin_right = Pt(4)
    ltf.margin_top = ltf.margin_bottom = Pt(3)
    lp = ltf.paragraphs[0]
    lp.text = legal_text or "—"
    _pptx_set_rtl_paragraph(lp, bold=False, size_pt=font_pt, color=_pptx_rgb("#000000"))


def _pptx_add_fields_table(
    slide,
    fields: list[dict[str, str]],
    *,
    top_in: float,
    height_in: float,
    margin_x: float,
    content_w: float,
) -> None:
    from pptx.util import Inches

    n = max(len(fields), 1)
    value_pt = 8 if n > 14 else (9 if n > 11 else 10)
    rows = n + 1
    table_shape = slide.shapes.add_table(
        rows, 2, Inches(margin_x), Inches(top_in), Inches(content_w), Inches(height_in)
    )
    table = table_shape.table
    table.columns[0].width = Inches(3.4)
    table.columns[1].width = Inches(9.03)
    _pptx_fill_cell_compact(table.cell(0, 0), "الحقل", header=True, size_pt=11)
    _pptx_fill_cell_compact(table.cell(0, 1), "القيمة", header=True, size_pt=11)
    if fields:
        for r, field in enumerate(fields, start=1):
            _pptx_fill_cell_compact(table.cell(r, 0), field["label"], size_pt=value_pt)
            _pptx_fill_cell_compact(table.cell(r, 1), field["value"], value=True, size_pt=value_pt)
    else:
        _pptx_fill_cell_compact(table.cell(1, 0), "—", size_pt=value_pt)
        _pptx_fill_cell_compact(table.cell(1, 1), "—", value=True, size_pt=value_pt)


def _pptx_add_footer(slide, margin_x: float, content_w: float, slide_h_in: float, exported_at: str) -> float:
    from pptx.util import Inches

    footer_h = 0.28
    foot_top = slide_h_in - footer_h - 0.05
    foot = slide.shapes.add_textbox(Inches(margin_x), Inches(foot_top), Inches(content_w), Inches(footer_h))
    fp = foot.text_frame.paragraphs[0]
    fp.text = f"تاريخ التصدير: {exported_at}"
    _pptx_set_rtl_paragraph(fp, size_pt=9, color=_pptx_rgb("#6b7280"))
    return foot_top


def _pptx_add_image_strip(
    slide,
    image_paths: list[Path],
    *,
    top_in: float,
    margin_x: float,
    content_w: float,
    strip_h: float,
) -> None:
    from pptx.util import Inches

    if not image_paths:
        return
    cap = slide.shapes.add_textbox(Inches(margin_x), Inches(top_in), Inches(content_w), Inches(0.2))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "الصور المرفقة"
    _pptx_set_rtl_paragraph(cp, bold=True, size_pt=9, color=_pptx_rgb("#14532d"))
    shown = image_paths[:3]
    slot_w_in = content_w / len(shown)
    for i, img_path in enumerate(shown):
        left_in = margin_x + slot_w_in * i + 0.05
        try:
            slide.shapes.add_picture(
                str(img_path),
                Inches(left_in),
                Inches(top_in + 0.22),
                width=Inches(slot_w_in - 0.1),
                height=Inches(max(0.5, strip_h - 0.28)),
            )
        except Exception:
            pass


def _build_legal_text_pptx(record: dict[str, object]) -> BytesIO:
    """Export legal text + fields. Long legal text gets a dedicated first slide (full text visible)."""
    from pptx import Presentation
    from pptx.util import Inches

    legal_text = str(record.get("legal_text") or "")
    fields = _pptx_export_fields(list(record.get("fields") or []))
    images = list(record.get("images") or [])
    exported_at = pd.Timestamp.now().strftime("%Y-%m-%d")

    image_paths: list[Path] = []
    for url in images:
        p = _image_url_to_path(str(url))
        if p is not None:
            image_paths.append(p)
    has_images = bool(image_paths)

    prs = Presentation()
    slide_w_in = 13.333
    slide_h_in = 7.5
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    blank = prs.slide_layouts[6]

    margin_x = 0.45
    content_w = slide_w_in - margin_x * 2
    footer_h = 0.28
    images_h = 1.15 if has_images else 0.0
    long_text = len(legal_text.strip()) > 150

    if long_text:
        # Slide 1 — full legal text only (never truncated)
        slide_text = prs.slides.add_slide(blank)
        bar_h = _pptx_add_header_bar(slide_text, prs, margin_x, content_w, "النص النظامي")
        text_top = bar_h + 0.1
        text_h = slide_h_in - text_top - footer_h - 0.15
        _pptx_add_legal_text_block(
            slide_text,
            legal_text,
            top_in=text_top,
            height_in=text_h,
            margin_x=margin_x,
            content_w=content_w,
            font_pt=11,
        )
        _pptx_add_footer(slide_text, margin_x, content_w, slide_h_in, exported_at)

        # Slide 2 — all detail fields (+ images)
        slide_fields = prs.slides.add_slide(blank)
        bar_h2 = _pptx_add_header_bar(slide_fields, prs, margin_x, content_w, "تفاصيل السجل")
        foot_top = _pptx_add_footer(slide_fields, margin_x, content_w, slide_h_in, exported_at)
        table_top = bar_h2 + 0.12
        table_h = max(1.0, foot_top - table_top - images_h - 0.12)
        _pptx_add_fields_table(
            slide_fields,
            fields,
            top_in=table_top,
            height_in=table_h,
            margin_x=margin_x,
            content_w=content_w,
        )
        if has_images:
            _pptx_add_image_strip(
                slide_fields,
                image_paths,
                top_in=table_top + table_h + 0.06,
                margin_x=margin_x,
                content_w=content_w,
                strip_h=images_h,
            )
    else:
        # Single slide — short legal text + fields
        slide = prs.slides.add_slide(blank)
        bar_h = _pptx_add_header_bar(slide, prs, margin_x, content_w, "النص النظامي")
        n_fields = len(fields)
        table_min_h = 0.32 * (max(n_fields, 1) + 1) + 0.3
        max_legal_h = max(0.7, slide_h_in - bar_h - table_min_h - images_h - footer_h - 0.4)
        legal_h_in, legal_pt = _pptx_legal_text_layout(legal_text, content_w, max_legal_h)
        legal_top = bar_h + 0.08
        _pptx_add_legal_text_block(
            slide,
            legal_text,
            top_in=legal_top,
            height_in=legal_h_in,
            margin_x=margin_x,
            content_w=content_w,
            font_pt=legal_pt,
        )
        foot_top = _pptx_add_footer(slide, margin_x, content_w, slide_h_in, exported_at)
        table_top = legal_top + legal_h_in + 0.1
        table_h = max(0.85, foot_top - table_top - images_h - 0.12)
        _pptx_add_fields_table(
            slide,
            fields,
            top_in=table_top,
            height_in=table_h,
            margin_x=margin_x,
            content_w=content_w,
        )
        if has_images:
            _pptx_add_image_strip(
                slide,
                image_paths,
                top_in=table_top + table_h + 0.06,
                margin_x=margin_x,
                content_w=content_w,
                strip_h=images_h,
            )

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _safe_pptx_filename(legal_text: str) -> str:
    base = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "", legal_text).strip() or "legal-text"
    if len(base) > 60:
        base = base[:60].rstrip()
    stamp = pd.Timestamp.now().strftime("%Y-%m-%d")
    return f"legal-text-{base}-{stamp}.pptx"


def _multi_filter_param_map() -> dict[str, str]:
    """عمود الفلترة في الإطار ← اسم معامل الاستعلام (يمكن تكراره: OR)."""
    return {
        COL_INHERENT_RISK: "inherent",
        COL_RESIDUAL_RISK: "residual",
        COL_STATUS: "status",
        YEAR_COL: "year",
        COL_DEPARTMENT: "department",
        COL_LEGISLATOR: "legislator",
        COL_SYSTEM_NAME: "system_name",
        COL_AUTHORITY: "authority",
        COL_REGULATION: "regulation",
        COL_LEGAL_TEXT: "legal_text",
        COL_COMPLIANCE_STATUS: "compliance_status",
        COL_CONTROL_CATEGORY: "control_category",
        COL_HOLDING_COMPANY: "holding_company",
        COL_SUBSIDIARY_COMPANY: "subsidiary_company",
    }


def _empty_multi_filters() -> dict[str, list[str]]:
    return {col: [] for col in _multi_filter_param_map()}


def _apply_filters(df: pd.DataFrame, selected: dict[str, list[str]], skip_key: str | None = None) -> pd.DataFrame:
    result = df
    for key, values in selected.items():
        if key == skip_key or not values:
            continue
        if key not in result.columns:
            continue
        result = result[result[key].isin(values)]
    return result


def _sort_values(key: str, values: list[str]) -> list[str]:
    if key == YEAR_COL:
        numeric = sorted([v for v in values if v.isdigit()], key=int)
        rest = sorted([v for v in values if not v.isdigit() and v != BLANK_LABEL])
        if BLANK_LABEL in values:
            return [BLANK_LABEL, *numeric, *rest]
        return [*numeric, *rest]
    # Mixed Excel cell types (e.g., int + str) can coexist in one column.
    # Sort by string representation to avoid TypeError on Python 3.
    return sorted(values, key=lambda v: str(v))


def _build_group_data(df: pd.DataFrame, key: str) -> list[dict[str, int | str]]:
    counts = df[key].value_counts(dropna=False).to_dict()
    sorted_keys = _sort_values(key, list(counts.keys()))
    return [{"key": k, "label": k, "count": int(counts[k])} for k in sorted_keys]
